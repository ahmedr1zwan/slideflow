import hashlib
from flask import Flask, Blueprint, request, jsonify, current_app
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
from pptx import Presentation
from google.cloud import vision
from fuzzywuzzy import fuzz  # Add fuzzy matching support
from sentence_transformers import SentenceTransformer, util
from utils import validate_file_path, allowed_file
import redis
import json


pptx_routes = Blueprint('pptx_routes', __name__)

# Initialize the Flask app
app = Flask(__name__)
CORS(app)  # Enable CORS

# Initialize Redis client
redis_client = redis.StrictRedis(host='localhost', port=6379, db=0, decode_responses=True)

# Initialize Google Vision client
vision_client = vision.ImageAnnotatorClient()

# Load Sentence-BERT model
sbert_model = SentenceTransformer('all-MiniLM-L6-v2')

@pptx_routes.route('/upload', methods=['POST'])
def upload_pptx():
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    if not allowed_file(file.filename, {'pptx'}):
        return jsonify({"error": "Unsupported file type. Only PPTX files are allowed."}), 400
    
    filename = secure_filename(file.filename)
    upload_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
    os.makedirs(os.path.dirname(upload_path), exist_ok=True)  # Create folder if not exists
    file.save(upload_path)
    return jsonify({"message": "File uploaded successfully", "path": upload_path}), 200

@pptx_routes.route('/analyze', methods=['POST'])
def analyze_pptx_content():
    data = request.get_json()
    if not data or 'file_path' not in data:
        return jsonify({"error": "File path not provided"}), 400

    file_path = data['file_path']
    error = validate_file_path(file_path)
    if error:
        return error

    # Generate a unique hash for the file
    file_hash = hashlib.md5(file_path.encode()).hexdigest()

    try:
        # Check if the file is already cached in Redis
        cached_data = redis_client.get(file_hash)
        if cached_data:
            return jsonify({"message": "Data retrieved from Redis cache", "data": eval(cached_data)}), 200

        # Extract and analyze content from the PowerPoint file
        presentation = Presentation(file_path)
        slides_data = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = ""
            image_analysis = []

            # Extract text and analyze images
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text + " "

                # Analyze images using Google Vision
                if shape.shape_type == 13:  # Shape type 13 corresponds to pictures
                    image = shape.image
                    image_format = image.ext
                    image_filename = f"slide_{slide_number}_image_{len(image_analysis) + 1}.{image_format}"
                    image_path = os.path.join(current_app.config['UPLOAD_FOLDER'], image_filename)

                    # Save the image
                    with open(image_path, "wb") as img_file:
                        img_file.write(image.blob)

                    # Perform analysis using Vision API
                    with open(image_path, "rb") as img_file:
                        content = img_file.read()
                        image = vision.Image(content=content)
                        response = vision_client.annotate_image({
                            "image": image,
                            "features": [
                                {"type": vision.Feature.Type.TEXT_DETECTION},
                                {"type": vision.Feature.Type.LABEL_DETECTION},
                                {"type": vision.Feature.Type.OBJECT_LOCALIZATION},
                                {"type": vision.Feature.Type.SAFE_SEARCH_DETECTION},
                                {"type": vision.Feature.Type.LOGO_DETECTION}
                            ],
                        })

                    # Extract analysis results
                    text = response.text_annotations[0].description if response.text_annotations else ""
                    labels = [label.description for label in response.label_annotations]
                    objects = [
                        {
                            "name": obj.name,
                            "bounding_poly": [
                                {"x": vertex.x, "y": vertex.y}
                                for vertex in obj.bounding_poly.normalized_vertices
                            ]
                        }
                        for obj in response.localized_object_annotations
                    ]
                    safe_search = {
                        "adult": response.safe_search_annotation.adult,
                        "violence": response.safe_search_annotation.violence,
                        "racy": response.safe_search_annotation.racy,
                    }
                    logos = [logo.description for logo in response.logo_annotations]

                    image_analysis.append({
                        "text": text,
                        "labels": labels,
                        "objects": objects,
                        "safe_search": safe_search,
                        "logos": logos
                    })

            # Derive slide name
            slide_name = slide_text.splitlines()[0] if slide_text.strip() else f"Slide {slide_number}"

            slides_data.append({
                "slide_number": slide_number,
                "slide_name": slide_name,
                "slide_text": slide_text.strip(),
                "image_analysis": image_analysis
            })

        # Cache the results in Redis
        redis_client.set(file_hash, json.dumps({
            "file_name": os.path.basename(file_path),
            "slides": slides_data
        }))

        return jsonify({"message": "Content analyzed and cached", "data": slides_data}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to process the PowerPoint file: {str(e)}"}), 500



@pptx_routes.route('/search', methods=['POST'])
def search_pptx_phrase():
    search_all = request.args.get('searchAll', 'false').lower() == 'true'  # Check if searchAll is true
    data = request.get_json()

    if not data or 'phrase' not in data:
        return jsonify({"error": "Phrase not provided"}), 400

    phrase = data['phrase']
    phrase_embedding = sbert_model.encode(phrase)  # Encode the phrase for semantic matching
    semantic_threshold = 0.6  # Minimum threshold for semantic similarity
    fuzzy_threshold = 80  # Minimum threshold for fuzzy matching
    matching_results = []

    try:
        if search_all:
            # Retrieve all cached keys from Redis
            keys = redis_client.keys()
        else:
            # Search within a specific file
            if 'file_path' not in data:
                return jsonify({"error": "file_path must be provided if searchAll is false"}), 400

            file_path = data['file_path']
            file_hash = hashlib.md5(file_path.encode()).hexdigest()

            if not redis_client.exists(file_hash):
                return jsonify({"error": f"File '{file_path}' not found in Redis cache. Please analyze it first."}), 404

            keys = [file_hash]

        # Perform the search
        for key in keys:
            # Retrieve cached file data from Redis
            cached_data = redis_client.get(key)
            if not cached_data:
                continue
            file_data = json.loads(cached_data)  # Parse JSON string into Python dictionary

            file_name = file_data['file_name']
            slides_data = file_data['slides']

            for slide in slides_data:
                slide_number = slide['slide_number']
                slide_name = slide['slide_name']
                slide_text = slide['slide_text']
                image_analysis = slide['image_analysis']

                # Syntactic similarity (fuzzy matching)
                fuzzy_score = fuzz.partial_ratio(phrase.lower(), slide_text.lower())

                # Semantic similarity (cosine similarity with embeddings)
                slide_embedding = sbert_model.encode(slide_text)
                semantic_similarity = util.cos_sim(phrase_embedding, slide_embedding).item()

                # Combine scores (weighted average)
                combined_score = 0.7 * semantic_similarity + 0.3 * (fuzzy_score / 100)

                if combined_score >= semantic_threshold:
                    matching_results.append({
                        "file_name": file_name,
                        "slide_name": slide_name,
                        "slide_number": slide_number,
                        "match_type": "text",
                        "syntactic_score": fuzzy_score,
                        "semantic_similarity": semantic_similarity,
                        "combined_score": combined_score
                    })
                    continue

                # Check image metadata
                for image in image_analysis:
                    labels = image['labels']
                    image_text = image['text']

                    # Syntactic and semantic similarity in image text
                    if image_text.strip():
                        image_text_embedding = sbert_model.encode(image_text)
                        image_semantic_similarity = util.cos_sim(phrase_embedding, image_text_embedding).item()
                        image_fuzzy_score = fuzz.partial_ratio(phrase.lower(), image_text.lower())
                        image_combined_score = 0.7 * image_semantic_similarity + 0.3 * (image_fuzzy_score / 100)

                        if image_combined_score >= semantic_threshold:
                            matching_results.append({
                                "file_name": file_name,
                                "slide_name": slide_name,
                                "slide_number": slide_number,
                                "match_type": "image_text",
                                "syntactic_score": image_fuzzy_score,
                                "semantic_similarity": image_semantic_similarity,
                                "combined_score": image_combined_score
                            })
                            break
                    # Syntactic and semantic similarity in image labels
                    for label in labels:
                        label_embedding = sbert_model.encode(label)
                        label_semantic_similarity = util.cos_sim(phrase_embedding, label_embedding).item()
                        label_fuzzy_score = fuzz.partial_ratio(phrase.lower(), label.lower())
                        label_combined_score = 0.7 * label_semantic_similarity + 0.3 * (label_fuzzy_score / 100)

                        if label_combined_score >= semantic_threshold:
                            matching_results.append({
                                "file_name": file_name,
                                "slide_name": slide_name,
                                "slide_number": slide_number,
                                "match_type": "image_label",
                                "syntactic_score": label_fuzzy_score,
                                "semantic_similarity": label_semantic_similarity,
                                "combined_score": label_combined_score
                            })
                            break

        # Sort results by combined score in descending order
        matching_results.sort(key=lambda x: x['combined_score'], reverse=True)

        return jsonify({
            "message": f"Phrase '{phrase}' search results:",
            "results": matching_results
        }), 200

    except Exception as e:
        return jsonify({"error": f"An error occurred during search: {str(e)}"}), 500
