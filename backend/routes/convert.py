import os
from flask import Blueprint, request, jsonify, current_app, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from PIL import Image
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# Define blueprint
pptx_to_pdf_blueprint = Blueprint('pptx_to_pdf', __name__)

# Allowed extensions
ALLOWED_EXTENSIONS = {'pptx'}

# Helper function: Check if file type is allowed
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Endpoint to upload and convert PPTX to PDF
@pptx_to_pdf_blueprint.route('/convert', methods=['POST'])
def convert_pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Unsupported file type. Only .pptx files are allowed."}), 400

    # Secure the filename and save the uploaded file
    filename = secure_filename(file.filename)
    upload_folder = current_app.config.get('UPLOAD_FOLDER', './uploads')
    converted_folder = current_app.config.get('CONVERTED_FOLDER', './converted')

    os.makedirs(upload_folder, exist_ok=True)
    os.makedirs(converted_folder, exist_ok=True)

    upload_path = os.path.join(upload_folder, filename)
    file.save(upload_path)

    # Convert PPTX to PDF
    try:
        pdf_filename = f"{os.path.splitext(filename)[0]}.pdf"
        pdf_path = os.path.join(converted_folder, pdf_filename)

        presentation = Presentation(upload_path)
        images = []
        temp_folder = os.path.join(converted_folder, "temp_images")
        os.makedirs(temp_folder, exist_ok=True)

        # Extract slides as images
        for i, slide in enumerate(presentation.slides):
            image_path = os.path.join(temp_folder, f"slide_{i + 1}.png")
            slide.shapes._spTree.blob.save(image_path)  # Save slide as an image
            images.append(image_path)

        # Create a PDF from extracted images
        c = canvas.Canvas(pdf_path, pagesize=letter)
        for image_path in images:
            image = Image.open(image_path)
            width, height = image.size
            aspect_ratio = width / height

            pdf_width, pdf_height = letter
            new_width = pdf_width
            new_height = new_width / aspect_ratio

            if new_height > pdf_height:
                new_height = pdf_height
                new_width = new_height * aspect_ratio

            x_offset = (pdf_width - new_width) / 2
            y_offset = (pdf_height - new_height) / 2

            c.drawImage(image_path, x_offset, y_offset, width=new_width, height=new_height)
            c.showPage()
        c.save()

        # Cleanup temporary images
        for image in images:
            os.remove(image)
        os.rmdir(temp_folder)

        return send_file(pdf_path, as_attachment=True, mimetype='application/pdf')

    except Exception as e:
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500
